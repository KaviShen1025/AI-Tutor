import os
from typing import Dict, Any, List
import PyPDF2
import docx
from pptx import Presentation
import re
from fastapi import HTTPException

class DocumentProcessingService:
    """Service for processing uploaded documents and extracting their content."""
    
    def __init__(self):
        """Initialize the document processing service."""
        pass
    
    async def process_document(self, file_path: str) -> Dict[str, Any]:
        """
        Process a document file and extract its content.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary containing the extracted content and metadata
        """
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail=f"File not found: {file_path}")
        
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension == '.pdf':
                return await self._process_pdf(file_path)
            elif file_extension == '.docx':
                return await self._process_docx(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                return await self._process_pptx(file_path)
            else:
                raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_extension}")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error processing document: {str(e)}")
    
    async def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Process a PDF document and extract its content."""
        content = []
        structure = []
        
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                
                # Extract metadata if available
                metadata = reader.metadata
                
                # Process each page
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        content.append(page_text)
                        
                        # Try to identify headings for structure
                        lines = page_text.split('\n')
                        for line in lines:
                            line = line.strip()
                            # Heuristic for headings: all caps, numbers at start, short length
                            if (line.isupper() and len(line) < 100) or re.match(r'^[0-9]+\..*', line):
                                structure.append({
                                    "level": 1,
                                    "heading": line,
                                    "page": i + 1
                                })
                
                return {
                    "content": "\n\n".join(content),
                    "structure": structure,
                    "metadata": {
                        "title": metadata.title if metadata.title else os.path.basename(file_path),
                        "author": metadata.author if metadata.author else "Unknown",
                        "page_count": len(reader.pages),
                        "content_type": "application/pdf"
                    },
                    "file_path": file_path
                }
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")
    
    async def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """Process a Word document and extract its content."""
        try:
            doc = docx.Document(file_path)
            
            content = []
            structure = []
            current_heading_level = 0
            
            # Extract content and structure
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)
                    
                    # Check for headings
                    if para.style.name.startswith('Heading'):
                        try:
                            level = int(para.style.name.split()[-1])
                        except ValueError:
                            level = current_heading_level
                            
                        structure.append({
                            "level": level,
                            "heading": para.text
                        })
                        current_heading_level = level
            
            # Extract document properties
            core_props = doc.core_properties
            
            return {
                "content": "\n\n".join(content),
                "structure": structure,
                "metadata": {
                    "title": core_props.title if core_props.title else os.path.basename(file_path),
                    "author": core_props.author if core_props.author else "Unknown",
                    "created": str(core_props.created) if core_props.created else None,
                    "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                },
                "file_path": file_path
            }
        except Exception as e:
            raise Exception(f"Error processing DOCX: {str(e)}")
    
    async def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Process a PowerPoint presentation and extract its content."""
        try:
            prs = Presentation(file_path)
            
            content = []
            structure = []
            
            # Process each slide
            for i, slide in enumerate(prs.slides):
                slide_content = []
                
                # Get slide title if available
                if slide.shapes.title and slide.shapes.title.text:
                    slide_title = slide.shapes.title.text
                    slide_content.append(slide_title)
                    
                    # Add to structure
                    structure.append({
                        "level": 1,
                        "heading": slide_title,
                        "slide": i + 1
                    })
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content.append(shape.text)
                
                content.append("\n".join(slide_content))
            
            return {
                "content": "\n\n".join(content),
                "structure": structure,
                "metadata": {
                    "title": os.path.basename(file_path),
                    "slide_count": len(prs.slides),
                    "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                },
                "file_path": file_path
            }
        except Exception as e:
            raise Exception(f"Error processing PPTX: {str(e)}")
